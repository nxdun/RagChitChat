"""
Document processing classes for handling PDF and PPTX files
"""
from abc import ABC, abstractmethod
import os
from pathlib import Path
import logging
from typing import List, Dict, Any, Optional

import pypdf
from pptx import Presentation

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor(ABC):
    """Base class for document processors"""
    
    @abstractmethod
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """Process a document and return chunks of text with metadata"""
        pass
    
    def save_chunks(self, chunks: List[Dict[str, Any]], output_dir: str, filename: str) -> None:
        """Save processed chunks to a text file"""
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            for chunk in chunks:
                f.write(f"--- Chunk from {chunk['metadata']['source']} "
                       f"(Page/Slide {chunk['metadata'].get('page_num', 'N/A')}) ---\n")
                f.write(chunk['content'])
                f.write("\n\n")
        
        logger.info(f"Saved processed content to {output_path}")


class PDFProcessor(DocumentProcessor):
    """Processor for PDF documents"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        logger.info(f"Processing PDF: {file_path}")
        chunks = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf = pypdf.PdfReader(file)
                
                # 💡 Extract text from each page and divide into overlapping chunks to maintain context
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text.strip():
                        # Simple chunking by character count
                        if len(text) <= self.chunk_size:
                            chunks.append({
                                'content': text,
                                'metadata': {
                                    'source': os.path.basename(file_path),
                                    'page_num': i + 1
                                }
                            })
                        else:
                            # Overlap chunking
                            start = 0
                            while start < len(text):
                                end = min(start + self.chunk_size, len(text))
                                chunk_text = text[start:end]
                                chunks.append({
                                    'content': chunk_text,
                                    'metadata': {
                                        'source': os.path.basename(file_path),
                                        'page_num': i + 1,
                                        'chunk_part': f"{start}-{end}"
                                    }
                                })
                                start += self.chunk_size - self.chunk_overlap
                                if start >= len(text):
                                    break
            
            logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            return []


class PPTXProcessor(DocumentProcessor):
    """Processor for PowerPoint documents"""
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        logger.info(f"Processing PPTX: {file_path}")
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            # Note: PowerPoint slides can have nested content and complex layouts
            # This simple extraction works for basic text, but might miss complex elements
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from all text-containing shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    chunks.append({
                        'content': '\n'.join(slide_text),
                        'metadata': {
                            'source': os.path.basename(file_path),
                            'page_num': i + 1
                        }
                    })
            
            logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            return []


class DocumentProcessorFactory:
    """Factory for creating document processors based on file extension"""
    
    @staticmethod
    def get_processor(file_path: str) -> Optional[DocumentProcessor]:
        """Get appropriate document processor based on file extension"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return PDFProcessor()
        elif ext == '.pptx':
            return PPTXProcessor()
        else:
            logger.warning(f"Unsupported file format: {ext}")
            return None
